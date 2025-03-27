from pptx import Presentation
from pptx.util import Inches
import xml.etree.ElementTree as ET
import os
from bs4 import BeautifulSoup

# ⚙️ Conversion px ➜ pouces relative à la taille réelle de la slide
def relative_px_to_inches(px, axis='x', slide_width_px=1150, slide_height_px=700, slide_inches=(11.98, 7.29)):
    if axis == 'x':
        return (float(px) / slide_width_px) * slide_inches[0]
    else:
        return (float(px) / slide_height_px) * slide_inches[1]

def extract_title(node):
    meta = node.find("metadata")
    if meta is not None:
        title = meta.find("title")
        if title is not None and title.text:
            return title.text.strip()
    return "Sans titre"

def parse_xml_to_slides(xml_file, media_dir):
    tree = ET.parse(xml_file)
    root = tree.getroot()
    course_structure = []

    for node in root.findall(".//node"):
        subnodes = node.findall(".//node")
        if subnodes:
            section_title = extract_title(node)
            section = {"section": section_title, "slides": []}
            for sub in subnodes:
                section["slides"].append(extract_slide_data(sub, media_dir))
            course_structure.append(section)
        else:
            slide = extract_slide_data(node, media_dir)
            course_structure.append({"section": None, "slides": [slide]})
    
    return course_structure

def extract_slide_data(node, media_dir):
    slide_data = {'title': extract_title(node), 'texts': [], 'images': []}
    screen = node.find(".//screen")
    if screen is None:
        return slide_data

    for img in screen.findall("image"):
        content = img.find("content")
        if content is not None and "file" in content.attrib:
            file_path = content.attrib['file'].replace("@/", "")
            design = img.find("design")
            if design is not None:
                slide_data['images'].append({
                    "file": os.path.join(media_dir, file_path),
                    "left": float(design.attrib.get("left", 0)),
                    "top": float(design.attrib.get("top", 0)),
                    "width": float(design.attrib.get("width", 1)),
                    "height": float(design.attrib.get("height", 1))
                })

    for txt in screen.findall("text"):
        content = txt.find("content")
        if content is not None:
            raw_html = content.text or ''
            soup = BeautifulSoup(raw_html, "html.parser")
            text = soup.get_text().strip()
            design = txt.find("design")
            if design is not None:
                slide_data['texts'].append({
                    "text": text,
                    "left": float(design.attrib.get("left", 0)),
                    "top": float(design.attrib.get("top", 0)),
                    "width": float(design.attrib.get("width", 5)),
                    "height": float(design.attrib.get("height", 1))
                })

    return slide_data

def generate_pptx(structure, output_path):
    prs = Presentation()

    slide_inches = (11.98, 7.29)
    slide_width_px, slide_height_px = 1150, 700

    prs.slide_width = Inches(slide_inches[0])
    prs.slide_height = Inches(slide_inches[1])
    blank_layout = prs.slide_layouts[6]
    section_slide_map = []

    for group in structure:
        section_start_idx = len(prs.slides)

        for slide_data in group['slides']:
            slide = prs.slides.add_slide(blank_layout)
            slide.name = slide_data['title'] or "Slide"

            for img in slide_data['images']:
                if os.path.exists(img['file']):
                    slide.shapes.add_picture(
                        img['file'],
                        left=Inches(relative_px_to_inches(img['left'], 'x', slide_width_px, slide_height_px, slide_inches)),
                        top=Inches(relative_px_to_inches(img['top'], 'y', slide_width_px, slide_height_px, slide_inches)),
                        width=Inches(relative_px_to_inches(img['width'], 'x', slide_width_px, slide_height_px, slide_inches)),
                        height=Inches(relative_px_to_inches(img['height'], 'y', slide_width_px, slide_height_px, slide_inches))
                    )

            for txt in slide_data['texts']:
                textbox = slide.shapes.add_textbox(
                    Inches(relative_px_to_inches(txt['left'], 'x', slide_width_px, slide_height_px, slide_inches)),
                    Inches(relative_px_to_inches(txt['top'], 'y', slide_width_px, slide_height_px, slide_inches)),
                    Inches(relative_px_to_inches(txt['width'], 'x', slide_width_px, slide_height_px, slide_inches)),
                    Inches(relative_px_to_inches(txt['height'], 'y', slide_width_px, slide_height_px, slide_inches))
                )
                tf = textbox.text_frame
                tf.text = txt['text']

        if group['section']:
            section_slide_map.append((section_start_idx, group['section']))

    for idx, name in section_slide_map:
        prs.slides._sldIdLst.insert(idx, prs.slides._sldIdLst[idx])
        sec_tag = prs.slides._sldIdLst[idx]
        sec_tag.set("name", name)

    prs.save(output_path)
