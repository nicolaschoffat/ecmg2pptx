import os
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# === CONFIG ===
CANVAS_WIDTH = 1150
CANVAS_HEIGHT = 700
PPT_WIDTH = Inches(13)  # corresponds to 1150px
PPT_HEIGHT = Inches(7.91)  # corresponds to 700px
MEDIA_DIR = '.'  # Same folder as script

# === Helpers ===
def percent_to_inches(val, total):
    return Inches((float(val) / total) * (PPT_WIDTH.inches if total == CANVAS_WIDTH else PPT_HEIGHT.inches))

def extract_value(el, attr):
    return float(el.attrib.get(attr, "0"))

# === Load and Parse XML ===
tree = ET.parse('course.xml')
root = tree.getroot()

# Target first <node> and <screen>
screen = root.find(".//screen")

# === Setup Presentation ===
prs = Presentation()
prs.slide_width = PPT_WIDTH
prs.slide_height = PPT_HEIGHT
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# === Process Images ===
for image in screen.findall("image"):
    design = image.find("design")
    content = image.find("content")
    src = content.attrib["file"].replace("@/", "").strip()

    left = percent_to_inches(design.attrib["left"], CANVAS_WIDTH)
    top = percent_to_inches(design.attrib["top"], CANVAS_HEIGHT)
    width = percent_to_inches(design.attrib["width"], CANVAS_WIDTH)
    height = percent_to_inches(design.attrib["height"], CANVAS_HEIGHT)

    image_path = os.path.join(MEDIA_DIR, src)
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    else:
        print(f"⚠️ Image not found: {src}")

# === Process Text ===
for text in screen.findall("text"):
    design = text.find("design")
    content = text.find("content")
    html = content.text or ""
    
    left = percent_to_inches(design.attrib["left"], CANVAS_WIDTH)
    top = percent_to_inches(design.attrib["top"], CANVAS_HEIGHT)
    width = percent_to_inches(design.attrib["width"], CANVAS_WIDTH)
    height = percent_to_inches(design.attrib["height"], CANVAS_HEIGHT)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "DÉONTOLOGIE - LE DÉCRET DANS SES GRANDES LIGNES ET L'ARTICLE 2"
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Tahoma"
    run.font.color.rgb = RGBColor(0x13, 0xAB, 0xB5)

# === Optional: Process Audio (manual in PPT) ===
for sound in screen.findall("sound"):
    content = sound.find("content")
    audio_src = content.attrib.get("file", "").replace("@/", "").strip()
    audio_path = os.path.join(MEDIA_DIR, audio_src)
    if os.path.exists(audio_path):
        print(f"🔊 Add audio manually in PowerPoint: {audio_src}")
    else:
        print(f"⚠️ Audio not found: {audio_src}")

# === Save PPTX ===
prs.save("slide1.pptx")
print("✅ slide1.pptx generated successfully.")
