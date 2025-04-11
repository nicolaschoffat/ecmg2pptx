import streamlit as st
import zipfile
import tempfile
import os
from bs4 import BeautifulSoup
from xml.etree import ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from html.parser import HTMLParser
from PIL import Image

px_to_pt = {
    20: 15,
    25: 18,
    30: 22,
    35: 26,
    40: 30,
    45: 34,
    50: 38
}

st.set_page_config(page_title="ECMG to PowerPoint Converter")
st.title("\U0001F4E4 Convertisseur ECMG vers PowerPoint")

uploaded_file = st.file_uploader("Upload un module ECMG (zip SCORM)", type="zip")

class HTMLtoPPTX(HTMLParser):
    def __init__(self, text_frame, style=None):
        super().__init__()
        self.tf = text_frame
        self.p = text_frame.paragraphs[0]
        self.style = {"bold": False, "italic": False}
        self.default_style = style or {}
        self.run = self.p.add_run()
        self.apply_style()

    def handle_starttag(self, tag, attrs):
        attrs = dict(attrs)

        if tag == "b":
            self.style["bold"] = True
        elif tag == "i":
            self.style["italic"] = True
        elif tag == "font":
            if "face" in attrs:
                self.default_style["font"] = attrs["face"]
            if "color" in attrs:
                self.default_style["fontcolor"] = attrs["color"]
            if "size" in attrs:
                try:
                    px = int(attrs["size"])
                    self.default_style["fontsize"] = px
                except:
                    pass
        elif tag == "br":
            self.p = self.tf.add_paragraph()
            self.p.alignment = self.tf.paragraphs[0].alignment  # hérite de l'alignement du premier
            
        self.run = self.p.add_run()
        self.apply_style()

    def handle_endtag(self, tag):
        if tag == "b":
            self.style["bold"] = False
        elif tag == "i":
            self.style["italic"] = False
        elif tag == "font":
            self.default_style.pop("font", None)
            self.default_style.pop("fontcolor", None)
            self.default_style.pop("fontsize", None)
        self.run = self.p.add_run()
        self.apply_style()

    def handle_data(self, data):
        self.run.text += data

    def apply_style(self):
        font = self.run.font
        font.bold = self.style.get("bold", False)
        font.italic = self.style.get("italic", False)

        if "font" in self.default_style:
            font.name = self.default_style["font"]
        if "fontcolor" in self.default_style:
            color = self.default_style["fontcolor"].lstrip("#")
            if len(color) == 6:
                try:
                    font.color.rgb = RGBColor.from_string(color.upper())
                except ValueError:
                    pass
        if "fontsize" in self.default_style:
            try:
                px = int(self.default_style["fontsize"])
                pt = px_to_pt.get(px, int(px * 0.75))
                font.size = Pt(pt)
            except:
                pass

def from_course(val, axis):
    if axis == "y":
        corrected = float(val) + 10.917
        px = corrected / 152.838 * 700
    else:
        px = float(val) / 149.351 * 1150
    return px * 0.01043

def from_look(val):
    return float(val) * 0.01043

def add_content_items_to_notes(screen, slide, type_name, label_icon):
    content_el = screen.find(f".//content[@type='{type_name}']")
    if content_el is None:
        return

    notes = slide.notes_slide.notes_text_frame
    bullet_lines = [f"{label_icon} Vue {type_name} :"]

    if type_name == "Cards":
        cards_wrapper = content_el.find("cards")
        if cards_wrapper is not None:
            for card in cards_wrapper.findall("card"):
                face_el = card.find("face")
                back_el = card.find("back")

                face_raw = "".join(face_el.itertext()) if face_el is not None else ""
                back_raw = "".join(back_el.itertext()) if back_el is not None else ""

                face_text = BeautifulSoup(face_raw, "html.parser").get_text().strip()
                back_text = BeautifulSoup(back_raw, "html.parser").get_text().strip()

                bullet_lines.append("• Face :")
                bullet_lines.append(face_text)
                bullet_lines.append("• Back :")
                bullet_lines.append(back_text)
                bullet_lines.append("")

    elif type_name == "Carousel" or type_name == "Vista":
        items_el = content_el.find("items")
        if items_el is not None:
            for item in items_el.findall("item"):
                raw = "".join(item.itertext()).strip()
                soup = BeautifulSoup(raw, "html.parser")
                for b in soup.find_all("b"):
                    b.insert_before("**")
                    b.insert_after("**")
                for i_tag in soup.find_all("i"):
                    i_tag.insert_before("_")
                    i_tag.insert_after("_")
                text = soup.get_text(separator="\n").strip()
                bullet_lines.append(f"• {text}")

    if len(bullet_lines) > 1:
        notes.text += "\n\n" + "\n".join(bullet_lines)

        # Label visuel sur slide
        label_box = slide.shapes.add_textbox(Inches(9.4), Inches(0.2), Inches(2.4), Inches(0.6))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{label_icon} Cartes {type_name}"
        font = run.font
        font.name = "Arial"
        font.size = Pt(12)
        font.bold = True
        p.alignment = PP_ALIGN.RIGHT

def add_consigne_boxes(screen, slide, style_map):
    for el in screen.findall("consigne"):
        content_el = el.find("content")
        if content_el is None or not content_el.text:
            continue

        text_id = el.attrib.get("id") or el.attrib.get("author_id")
        style = style_map.get(text_id, {})
        design_el = el.find("design")

        def has_position_attrs(d):
            return (
                d is not None and any(
                    attr in d.attrib and float(d.attrib[attr]) > 0
                    for attr in ["top", "left", "width", "height"]
                )
            )

        if has_position_attrs(design_el):
            top_px = float(design_el.attrib.get("top", 0))
            left_px = float(design_el.attrib.get("left", 0))
            width_px = float(design_el.attrib.get("width", 140))
            height_px = float(design_el.attrib.get("height", 10))

            top = from_course(top_px, "y")
            left = from_course(left_px, "x")
            width = from_course(width_px, "x")
            height = from_course(height_px, "y")
        else:
            top_px = float(style.get("top", 0))
            left_px = float(style.get("left", 0))
            width_px = float(style.get("width", 140))
            height_px = float(style.get("height", 10))

            top = from_look(top_px)
            left = from_look(left_px)
            width = from_look(width_px)
            height = from_look(height_px)

        box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1), Inches(width), Inches(height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True

        alignment = style.get("align", "").lower()
        if alignment == "center":
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        elif alignment == "right":
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        else:
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        if "valign" in style:
            valign = style.get("valign", "").lower()
            if valign == "middle":
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            elif valign == "bottom":
                tf.vertical_anchor = MSO_ANCHOR.BOTTOM
            else:
                tf.vertical_anchor = MSO_ANCHOR.TOP

        parser = HTMLtoPPTX(tf, style)
        parser.feed(content_el.text)


# 🔧 Ajout fonction pour traiter les fichiers externes (PDF)
def add_external_links(screen, slide):
    for action in screen.iter("action"):
        if action.attrib.get("action") == "open":
            param = action.attrib.get("param", "")
            if param.endswith(".pdf") and param.startswith("@/"):
                # 1. Ajouter dans les notes
                notes = slide.notes_slide.notes_text_frame
                notes.text += f"\n\nLien vers un document externe : {param}"

                # 2. Ajouter un pictogramme textuel sur la slide (top-right par défaut)
                box = slide.shapes.add_textbox(Inches(10), Inches(0.3), Inches(2), Inches(0.5))
                tf = box.text_frame
                tf.word_wrap = True
                tf.text = "📎 Voir document joint"
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

if uploaded_file:
    with tempfile.TemporaryDirectory() as tmpdir:
        zip_path = os.path.join(tmpdir, "module.zip")
        with open(zip_path, "wb") as f:
            f.write(uploaded_file.read())
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(tmpdir)

        course_path, look_path, author_path = None, None, None
        for root, dirs, files in os.walk(tmpdir):
            if "course.xml" in files:
                course_path = os.path.join(root, "course.xml")
            if "look.xml" in files:
                look_path = os.path.join(root, "look.xml")
            if "author.xml" in files:
                author_path = os.path.join(root, "author.xml")

        if not course_path or not look_path or not author_path:
            st.error("Fichiers course.xml, look.xml ou author.xml introuvables.")
            st.stop()

        tree = ET.parse(course_path)
        root = tree.getroot()
        nodes = root.findall(".//node")

        # Lecture du titre global de l'UA (Unité d'Apprentissage)
        ua_title = "[Titre UA manquant]"
        metadata = root.find("./metadata")
        if metadata is not None:
            global_title_el = metadata.find("title")
            if global_title_el is not None and global_title_el.text:
                ua_title = global_title_el.text.strip()
        
        look_tree = ET.parse(look_path)
        look_root = look_tree.getroot()
        style_map = {}
        for el in look_root.findall(".//*[@id]"):
            design = el.find("design")
            if design is not None:
                style_map[el.attrib["id"]] = design.attrib
                if "author_id" in el.attrib:
                    style_map[el.attrib["author_id"]] = design.attrib

        author_tree = ET.parse(author_path)
        author_root = author_tree.getroot()
        author_map = {
            el.attrib.get("id"): el.findtext("description")
            for el in author_root.findall(".//item")
        }

        prs = Presentation()
        prs.slide_width = Inches(12)
        prs.slide_height = Inches(7.3)

        for i, node in enumerate(nodes):
            title_el = node.find("./metadata/title")
            title_text = title_el.text.strip() if title_el is not None and title_el.text else "Sans titre"
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # 🔁 Injecter des éléments de look.xml spécifiques à certaines pages
            look_elements_by_page = {
                "page_intro": ["cadre_intro", "title_UA_intro"]
            }
            
            page_id = node.attrib.get("id")
            if page_id in look_elements_by_page:
                for el_id in look_elements_by_page[page_id]:
                    look_el = look_root.find(f".//*[@id='{el_id}']")
                    if look_el is not None:
                        tag = look_el.tag
                        style = style_map.get(el_id, {})
                        design_el = look_el.find("design")
                        content_el = look_el.find("content")
            
                        def has_position_attrs(d):
                            return d is not None and any(attr in d.attrib and float(d.attrib[attr]) > 0 for attr in ["top", "left", "width", "height"])
            
                        if has_position_attrs(design_el):
                            top_px = float(design_el.attrib.get("top", 0))
                            left_px = float(design_el.attrib.get("left", 0))
                            width_px = float(design_el.attrib.get("width", 140))
                            height_px = float(design_el.attrib.get("height", 10))
                            top = from_look(top_px)
                            left = from_look(left_px)
                            width = from_look(width_px)
                            height = from_look(height_px)
            
                            if tag == "image" and content_el is not None and content_el.attrib.get("file"):
                                image_path = os.path.join(os.path.dirname(look_path), content_el.attrib["file"])
                                if os.path.exists(image_path):
                                    with Image.open(image_path) as im:
                                        orig_width_px, orig_height_px = im.size
                                    orig_ratio = orig_width_px / orig_height_px
                                    target_ratio = width / height
                                    if orig_ratio > target_ratio:
                                        draw_width = width
                                        draw_height = width / orig_ratio
                                        offset_left = 0
                                        offset_top = (height - draw_height) / 2
                                    else:
                                        draw_height = height
                                        draw_width = height * orig_ratio
                                        offset_top = 0
                                        offset_left = (width - draw_width) / 2
                                    slide.shapes.add_picture(
                                        image_path,
                                        Inches(left + offset_left + 0.1),
                                        Inches(top + offset_top + 0.1),
                                        width=Inches(draw_width),
                                        height=Inches(draw_height)
                                    )
                            elif tag in ["text", "title"]:
                                box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1), Inches(width), Inches(height))
                                tf = box.text_frame
                                tf.clear()
                                tf.word_wrap = True
                                parser = HTMLtoPPTX(tf, style)
                                if content_el is None or not content_el.text:
                                    parser.feed(ua_title)
                                else:
                                    parser.feed(content_el.text)
            
            title_style = style_map.get("titre_activite")
            title_shape = slide.shapes.title
            
            if title_style:
                try:
                    top = from_look(float(title_style.get("top", 0)))
                    left = from_look(float(title_style.get("left", 0)))
                    width = from_look(float(title_style.get("width", 800)))
                    height = from_look(float(title_style.get("height", 50)))
            
                    title_shape.left = Inches(left + 0.1)
                    title_shape.top = Inches(top + 0.1)
                    title_shape.width = Inches(width)
                    title_shape.height = Inches(height)
                except Exception as e:
                    st.warning(f"❗ Erreur redimension titre: {e}")
            
            tf = title_shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title_text
            
            font = run.font
            font.name = title_style.get("font", "Tahoma") if title_style else "Tahoma"
            try:
                fontsize = int(title_style.get("fontsize", 22)) if title_style else 22
                font.size = Pt(px_to_pt.get(fontsize, int(fontsize * 0.75)))
            except:
                font.size = Pt(16.5)
            
            font.bold = title_style.get("bold", "0") == "1" if title_style else False
            font.italic = title_style.get("italic", "0") == "1" if title_style else False
            
            color = title_style.get("fontcolor", "#000000").lstrip("#") if title_style else "000000"
            if len(color) == 6:
                try:
                    font.color.rgb = RGBColor.from_string(color.upper())
                except ValueError:
                    pass
            
            align = title_style.get("align", "left").lower() if title_style else "left"
            if align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif align == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            page = node.find(".//page")
            screen = page.find("screen") if page is not None else None
            if not screen:
                continue

            page_type = node.find("page").attrib.get("type", "") if node.find("page") is not None else ""

            if page_type == "result":
            # 📌 Label "Page Bilan"
            box = slide.shapes.add_textbox(Inches(9.4), Inches(0.2), Inches(2.4), Inches(0.6))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "📊 Page Bilan"
            font = run.font
            font.name = "Arial"
            font.size = Pt(12)
            font.bold = True
            p.alignment = PP_ALIGN.RIGHT
        
            # 📋 Notes avec les textes pour chaque score
            page_el = node.find("page")
            results_el = page_el.find("results") if page_el is not None else None
            if results_el is not None:
                notes = slide.notes_slide.notes_text_frame
                notes.text += "\n\n🧾 Résultats affichés selon score :"
        
                for result in results_el.findall("result"):
                    score = result.attrib.get("score", "?")
                    screen_result = result.find("screen")
                    text_block = screen_result.find("text") if screen_result is not None else None
                    content_el = text_block.find("content") if text_block is not None else None
        
                    raw = "".join(content_el.itertext()) if content_el is not None else ""
                    text_clean = BeautifulSoup(raw, "html.parser").get_text().strip()
        
                    notes.text += f"\n---\n🔢 Score {score} :\n{text_clean}"
        
            else:
                st.warning(f"Aucune balise <results> trouvée dans node id={node.attrib.get('id')}")
           
            # ✅ Ajout de contenu spécifique selon type (Vista, Cards, Carousel)
            add_content_items_to_notes(screen, slide, "Vista", "🪟")
            add_content_items_to_notes(screen, slide, "Cards", "🃏")
            add_content_items_to_notes(screen, slide, "Carousel", "🎠")
            
            # ✅ Ajout des consignes au début du traitement de l'écran
            add_consigne_boxes(screen, slide, style_map)
            
            # ✅ Ajout des liens vers documents PDF
            add_external_links(screen, slide)
            
            y = 1.5
            # 🎥 Vidéo (si présente)
            video_file = None
            for content_el in screen.findall(".//content"):
                if "file" in content_el.attrib and content_el.attrib["file"].endswith(".mp4"):
                    video_file = content_el.attrib["file"]
                    break
            if video_file:
                box = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(6), Inches(1))
                tf = box.text_frame
                tf.text = f" Vidéo : {video_file} à intégrer"
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 🎞️ Flash (animation à convertir)
            for flash_el in screen.findall(".//flash"):
                content_el = flash_el.find("content")
                if content_el is not None and "file" in content_el.attrib:
                    flash_file = content_el.attrib["file"]
                    # Ajout d'un pictogramme texte sur la slide
                    box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(10), Inches(0.6))
                    tf = box.text_frame
                    tf.word_wrap = True
                    tf.text = f"🎞️ Animation Flash à recréer ou convertir depuis ECMG : {flash_file}"
                    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            
                    # Ajout dans les notes pour l'auteur
                    notes = slide.notes_slide.notes_text_frame
                    notes.text += f"\nContenu Flash détecté : {flash_file}"
        
            # 🔊 Sons (audio)
            sound_blocks = screen.findall(".//sound")
            if sound_blocks:
                notes = slide.notes_slide.notes_text_frame
                audio_notes = []
                for snd in sound_blocks:
                    author_id = snd.attrib.get("author_id")
                    content = snd.find("content")
                    filename = content.attrib.get("file") if content is not None else None
                    audio_text = author_map.get(author_id)
                    if filename:
                        audio_notes.append(f"Audio : {filename}\nTexte lu : {audio_text or '[non trouvé]'}")
                if audio_notes:
                    notes.text += "\n\n" + "\n---\n".join(audio_notes)
            # ❓ QCM (MCQText)
            elfe = screen.find("elfe")
            if elfe is not None and elfe.find("content") is not None and elfe.find("content").attrib.get("type") == "MCQText":
                items = elfe.find("content/items")
                question_el = screen.find("question")
                if question_el is not None:
                    question_text = BeautifulSoup(question_el.find("content").text, "html.parser").get_text()
                    box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(10), Inches(1))
                    box.text_frame.text = f"❓ {question_text}"
                    y += 1.0
                for item in items.findall("item"):
                    score = item.attrib.get("score", "0")
                    label = "✅" if int(score) > 0 else "⬜"
                    box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(9.5), Inches(0.5))
                    box.text_frame.text = f"{label} {item.text.strip()}"
                    y += 0.5
                feedbacks = page.findall(".//feedbacks/correc/fb/screen/feedback")
                notes = slide.notes_slide.notes_text_frame
                feedback_texts = []
                for fb in feedbacks:
                    fb_content = fb.find("content")
                    if fb_content is not None and fb_content.text:
                        soup = BeautifulSoup(fb_content.text, "html.parser")
                        feedback_texts.append(soup.get_text(separator="\n"))
                if feedback_texts:
                    notes.text += "\n---\n" + "\n---\n".join(feedback_texts)
            # 🖼️ Images dans le screen
            # ✅ Gestion de la profondeur en suivant l'ordre d'apparition dans le XML
            for el in list(screen):
                tag = el.tag
            
                if tag == "image":
                    content = el.find("content")
                    if content is None or not content.attrib.get("file"):
                        continue
                    img_file = content.attrib["file"]
                    image_id = el.attrib.get("id") or el.attrib.get("author_id")
                    style = style_map.get(image_id, {})
                    design_el = el.find("design")
            
                    def has_position_attrs(d):
                        return (
                            d is not None and any(
                                attr in d.attrib and float(d.attrib[attr]) > 0
                                for attr in ["top", "left", "width", "height"]
                            )
                        )
            
                    if has_position_attrs(design_el):
                        top_px = float(design_el.attrib.get("top", 0))
                        left_px = float(design_el.attrib.get("left", 0))
                        width_px = float(design_el.attrib.get("width", 200))
                        height_px = float(design_el.attrib.get("height", 200))
                        top = from_course(top_px, "y")
                        left = from_course(left_px, "x")
                        width = from_course(width_px, "x")
                        height = from_course(height_px, "y")
                    else:
                        top_px = float(style.get("top", 0))
                        left_px = float(style.get("left", 0))
                        width_px = float(style.get("width", 200))
                        height_px = float(style.get("height", 200))
                        top = from_look(top_px)
                        left = from_look(left_px)
                        width = from_look(width_px)
                        height = from_look(height_px)
            
                    image_dir = os.path.dirname(course_path)
                    image_path = os.path.join(image_dir, os.path.basename(img_file))
                    if not os.path.exists(image_path):
                        image_dir = os.path.dirname(look_path)
                        image_path = os.path.join(image_dir, os.path.basename(img_file))
                    if os.path.exists(image_path):
                        try:
                            with Image.open(image_path) as im:
                                orig_width_px, orig_height_px = im.size
                            orig_ratio = orig_width_px / orig_height_px
                            target_ratio = width / height
            
                            if orig_ratio > target_ratio:
                                draw_width = width
                                draw_height = width / orig_ratio
                                offset_left = 0
                                offset_top = (height - draw_height) / 2
                            else:
                                draw_height = height
                                draw_width = height * orig_ratio
                                offset_top = 0
                                offset_left = (width - draw_width) / 2
            
                            slide.shapes.add_picture(
                                image_path,
                                Inches(left + offset_left + 0.1),
                                Inches(top + offset_top + 0.1),
                                width=Inches(draw_width),
                                height=Inches(draw_height)
                            )
                        except Exception as e:
                            st.warning(f"⚠️ Erreur ajout image {img_file} : {e}")
            
                elif tag == "text":
                    content_el = el.find("content")
                    if content_el is None or not content_el.text:
                        continue
            
                    text_id = el.attrib.get("id") or el.attrib.get("author_id")
                    style = style_map.get(text_id, {})
                    design_el = el.find("design")
            
                    def has_position_attrs(d):
                        return (
                            d is not None and any(
                                attr in d.attrib and float(d.attrib[attr]) > 0
                                for attr in ["top", "left", "width", "height"]
                            )
                        )
            
                    if has_position_attrs(design_el):
                        top_px = float(design_el.attrib.get("top", 0))
                        left_px = float(design_el.attrib.get("left", 0))
                        width_px = float(design_el.attrib.get("width", 140))
                        height_px = float(design_el.attrib.get("height", 10))
                        top = from_course(top_px, "y")
                        left = from_course(left_px, "x")
                        width = from_course(width_px, "x")
                        height = from_course(height_px, "y")
                    else:
                        top_px = float(style.get("top", 0))
                        left_px = float(style.get("left", 0))
                        width_px = float(style.get("width", 140))
                        height_px = float(style.get("height", 10))
                        top = from_look(top_px)
                        left = from_look(left_px)
                        width = from_look(width_px)
                        height = from_look(height_px)
            
                    box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1), Inches(width), Inches(height))
                    tf = box.text_frame
                    tf.clear()
                    tf.word_wrap = True
            
                    alignment = style.get("align", "").lower()
                    if alignment == "center":
                        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    elif alignment == "right":
                        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                    else:
                        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            
                    if "valign" in style:
                        valign = style.get("valign", "").lower()
                        if valign == "middle":
                            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                        elif valign == "bottom":
                            tf.vertical_anchor = MSO_ANCHOR.BOTTOM
                        else:
                            tf.vertical_anchor = MSO_ANCHOR.TOP
            
                    parser = HTMLtoPPTX(tf, style)
                    if content_el is None or not content_el.text:
                        parser.feed(ua_title)
                    else:
                        parser.feed(content_el.text)

        output_path = os.path.join(tmpdir, "converted.pptx")
        prs.save(output_path)

        with open(output_path, "rb") as f:
            st.download_button("📅 Télécharger le PowerPoint", data=f, file_name="module_ecmg_converti.pptx")
